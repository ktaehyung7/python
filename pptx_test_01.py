from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Cm
import openpyxl
import pandas as pd


prs = Presentation("C:/Users/ktaeh/Documents/감리/제안서/##분류##/재구축_개인정보/20230411_서울특별시_청년몽땅정보통_전면_재구축_감리_및_개인정보영향평가/제안서/서울시 청년몽땅정보통 전면 재구축 감리 및 개인정보영향평가 용역_제안서(원본)_v0.26.pptx")

wb = openpyxl.load_workbook('C:/Users/ktaeh/Documents/감리/제안서/##분류##/재구축_개인정보/20230411_서울특별시_청년몽땅정보통_전면_재구축_감리_및_개인정보영향평가/제안서/test01.xlsx')
ws = wb.get_sheet_by_name('이력')


id_lists = ['{e6_nm}', \
            '{e6_ca1}', '{e6_ca2}', '{e6_ca3}', \
            '{e6_ca4}']

res_lists = []
for row in ws.rows:
    res_lists.append(row[2].value)


for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if "{p23_img}" in run.text:
                        left=Cm(2.0);top=Cm(25.3);width=Cm(1.5);height=Cm(2.0)
                        pic=slide.shapes.add_picture('./김청열.png', left, top, width, height)
                    for id_list in id_lists:
                        if id_list in run.text:
                            run.text = run.text.replace(id_list, str(res_lists[id_lists.index(id_list)]))
        i=0
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for shp in shape.shapes:
                if shp.has_text_frame:
                    for paragraph in shp.text_frame.paragraphs:
                        for run in paragraph.runs:
                            for id_list in id_lists:
#                                if '{id_3_tca}' in run.text:
#                                    run.text = run.text.replace('{id_3_tca}', str('인사관리시스템 구축'))
                                if id_list in run.text:
                                    run.text = run.text.replace(id_list, str(res_lists[id_lists.index(id_list)]))
        i=0
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            for id_list in id_lists:
                                if id_list in run.text:
                                    run.text = run.text.replace(id_list, str(res_lists[id_lists.index(id_list)]))

prs.save('./test1.pptx')
