from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Cm
import openpyxl
import pandas as pd


prs = Presentation("./test.pptx")

wb = openpyxl.load_workbook('./test01.xlsx')
ws = wb.get_sheet_by_name('이력')


id_lists = ['{id_3_name}', '{id3_cnt}', '{id_3_year}', \
            '{id_3_lic}', '{id_3_li_cnt}', \
            '{tb_3_ca1}', '{tb_3_ca2}', '{tb_3_ca3}', \
            '{tb_3_ca4}', '{tb_3_ca5}', '{tb_3_ca6}', \
            '{tb_3_ca7}', '{tb_3_ca8}', '{tb_3_ca9}', \
            '{tb_3_ca10}', '{tb_3_ca11}', '{tb_3_ca12}']

res_lists = []
for row in ws.rows:
    res_lists.append(row[2].value)


for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if "{id_89_name3_img}" in run.text:
                        left=Cm(18.3);top=Cm(7.3);width=Cm(1.5);height=Cm(2.0)
                        pic=slide.shapes.add_picture('./김태형.jpg', left, top, width, height)
        i=0
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for shp in shape.shapes:
                if shp.has_text_frame:
                    for paragraph in shp.text_frame.paragraphs:
                        for run in paragraph.runs:
                            for id_list in id_lists:
                                if '{id_3_tca}' in run.text:
                                    run.text = run.text.replace('{id_3_tca}', str('인사관리시스템 구축'))
                                if id_list in run.text:
                                    run.text = run.text.replace(id_list, str(res_lists[id_lists.index(id_list)]))
        i=0
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            for id_list in id_lists:
                                if id_list in run.text:
                                    run.text = run.text.replace(id_list, str(res_lists[id_lists.index(id_list)]))

prs.save('./test1.pptx')
